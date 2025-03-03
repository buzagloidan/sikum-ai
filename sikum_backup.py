#!/usr/bin/env python3.10
import sys
import os
import inspect

# Add virtualenv site-packages to path
VIRTUALENV_PATH = '/usr/local/lib/python3.10/site-packages'
if VIRTUALENV_PATH not in sys.path:
    sys.path.append(VIRTUALENV_PATH)

import logging
import random
from telegram import (
    Update,
    InlineKeyboardButton,
    InlineKeyboardMarkup,
    Bot,
    LabeledPrice
)
from telegram.ext import (
    ApplicationBuilder,
    CommandHandler,
    CallbackQueryHandler,
    ContextTypes,
    MessageHandler,
    filters,
    PreCheckoutQueryHandler,
    ConversationHandler
)
import google.generativeai as genai
import PyPDF2
import json
import asyncio
from collections import defaultdict, Counter
import time
from telegram.constants import FileSizeLimit
import math
import sqlite3
from datetime import datetime, date, timedelta
from contextlib import contextmanager

# Import our new optimized database utilities
from db_utils import (
    get_db_connection, 
    get_async_db_connection,
    resilient_db_operation,
    query_cache,
    get_quiz_with_caching as db_get_quiz_with_caching,
    get_user_quizzes_with_caching as db_get_user_quizzes_with_caching,
    optimize_db,
    vacuum_db,
    init_db,
    start_maintenance_task
)

# NEW: import for reading DOCX files
from docx import Document as DocxDocument
from telegram.error import BadRequest, TimedOut, NetworkError, RetryAfter, TelegramError

# Add to imports
from pptx import Presentation  # For PPTX files
from typing import Optional, Dict, List, Union, Tuple, TypeVar, AsyncGenerator, Generator
import traceback
from functools import lru_cache, wraps
from concurrent.futures import ThreadPoolExecutor
from tenacity import retry, stop_after_attempt, wait_exponential, retry_if_exception_type

def with_network_retry(func):
    @retry(
        stop=stop_after_attempt(3),
        wait=wait_exponential(multiplier=1, min=4, max=10),
        retry=retry_if_exception_type((TimedOut, NetworkError, RetryAfter))
    )
    @wraps(func)
    async def wrapper(*args, **kwargs):
        try:
            return await func(*args, **kwargs)
        except (TimedOut, NetworkError) as e:
            logging.warning(f"Network error in {func.__name__}: {str(e)}, retrying...")
            raise
    return wrapper
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Replace hardcoded values with environment variables
BOT_TOKEN = os.getenv('BOT_TOKEN') or ''
if not BOT_TOKEN:
    raise ValueError("BOT_TOKEN environment variable not set")
GEMINI_API_KEY = os.getenv('GEMINI_API_KEY')
ADMIN_IDS = [int(id.strip()) for id in os.getenv('ADMIN_IDS', '').split(',')]
DB_PATH = os.getenv('DB_PATH', 'bot_stats.db')
PAYMENT_PROVIDER_TOKEN = os.getenv('PAYMENT_PROVIDER_TOKEN', '')
SUBSCRIPTION_COST = int(os.getenv('SUBSCRIPTION_COST', '100'))
FREE_USES_PER_DAY = int(os.getenv('FREE_USES_PER_DAY', '15'))  # Changed from '3' to '15'
PREMIUM_USES_PER_DAY = int(os.getenv('PREMIUM_USES_PER_DAY', '30'))
MAX_FILE_SIZE_MB = int(os.getenv('MAX_FILE_SIZE_MB', '50'))
MAX_PAGES_TOTAL = int(os.getenv('MAX_PAGES_TOTAL', '100'))
TEMP_DIR = os.getenv('TEMP_DIR', 'temp')
SUPPORT_CONTACT = os.getenv('SUPPORT_CONTACT', '@sikumai')
RATE_LIMIT_PERIOD = int(os.getenv('RATE_LIMIT_PERIOD', '300'))
MAX_REQUESTS_PER_PERIOD = int(os.getenv('MAX_REQUESTS_PER_PERIOD', '5'))
VIRTUALENV_PATH = os.getenv('VIRTUALENV_PATH', '/usr/local/lib/python3.10/site-packages')



from logging.handlers import RotatingFileHandler

# Create logger instance
logger = logging.getLogger(__name__)

# Then use logger
if not GEMINI_API_KEY:
    logger.error("GEMINI_API_KEY environment variable not set")
    raise ValueError("GEMINI_API_KEY environment variable not set")
# Configure Gemini API
genai.configure(api_key=GEMINI_API_KEY)
# Set up model with optimal performance parameters
generation_config = {
    "temperature": 0.3,  # Lower temperature for more consistent, faster responses
    "top_p": 0.85,       # Slightly reduce top_p for faster responses
    "top_k": 40,         # Slightly increase top_k for better quality
    "max_output_tokens": 2048  # Limit output size for faster responses
}
model = genai.GenerativeModel('gemini-2.0-flash', generation_config=generation_config)

# ====== Logging Configuration ======
# Create logs directory in the same folder as the script
LOG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logs")
os.makedirs(LOG_DIR, exist_ok=True)

def setup_logging():
    """Configure logging with proper format and handlers."""
    # Create logs directory
    os.makedirs(LOG_DIR, exist_ok=True)
    
    # Configure root logger
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(name)s - %(levelname)s - %(funcName)s:%(lineno)d - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S',
        handlers=[
            RotatingFileHandler(
                os.path.join(LOG_DIR, 'bot.log'),
                maxBytes=10_000_000,  # 10MB
                backupCount=5
            ),
            logging.StreamHandler(sys.stdout)
        ]
    )
    
    # Set specific log levels for noisy libraries
    logging.getLogger('telegram').setLevel(logging.WARNING)
    logging.getLogger('httpx').setLevel(logging.WARNING)

# Initialize logging
setup_logging()

# Set telegram logger to WARNING
logging.getLogger("telegram").setLevel(logging.WARNING)
logging.getLogger("apscheduler").setLevel(logging.WARNING)

# ====== Constants ======
STATE_QUIZ_ACTIVE = "quiz_active"
QUESTIONS_PER_QUIZ = 20
MAX_FILE_SIZE = MAX_FILE_SIZE_MB * 1024 * 1024
MAX_PAGES_TOTAL = 100  # Max pages for PDFs
RATE_LIMIT_PERIOD = 300  # 5 minutes
MAX_REQUESTS_PER_PERIOD = 5

# Track active users and user request times
active_users = set()
user_request_times = defaultdict(list)

# Temporary directory
os.makedirs(TEMP_DIR, exist_ok=True)

# User Statistics Tracking
user_stats = {
    "total_users": set(),  # Set of all unique user IDs
    "daily_users": defaultdict(set),  # Dictionary mapping dates to sets of user IDs
    "total_quizzes": 0,  # Total number of quizzes generated
    "user_quiz_counts": defaultdict(int),  # Dictionary mapping user IDs to their quiz count
}

# Add new constants for subscription tracking
DB_SUBSCRIPTIONS_TABLE = """
CREATE TABLE IF NOT EXISTS subscriptions (
    user_id INTEGER PRIMARY KEY,
    subscribed_until DATE,
    FOREIGN KEY (user_id) REFERENCES users (user_id)
)
"""

# Add this with other constants
DB_DAILY_USAGE_TABLE = """
CREATE TABLE IF NOT EXISTS daily_usage (
    user_id INTEGER,
    date DATE,
    attempts INTEGER DEFAULT 1,
    PRIMARY KEY (user_id, date)
)
"""

# Add with other constants
HEBREW_LETTERS = ['א', 'ב', 'ג', 'ד']
MAX_SAVED_QUIZZES_DISPLAY = 10  # For /list and /play commands

# Update allowed extensions from a set to a dictionary
ALLOWED_EXTENSIONS = {
    '.pdf': 'pdf',
    '.docx': 'docx', 
    '.doc': 'docx',    # Old Word format
    '.txt': 'txt',     # Plain text
    '.pptx': 'pptx',   # PowerPoint
    '.ppt': 'pptx'     # Old PowerPoint format
}

# ====== Custom Exceptions ======
class FileProcessingError(Exception):
    """Raised when file processing fails"""
    pass

class BotError(Exception):
    """Base exception for bot errors."""
    pass

class DatabaseError(BotError):
    """Database related errors."""
    pass

class RateLimitError(BotError):
    """Rate limiting errors."""
    pass

# ====== File Processing Classes ======
class FileProcessor:
    """Handles file processing pipeline."""
    
    def __init__(self, file_path: str, file_type: str):
        self.file_path = file_path
        self.file_type = file_type
        
    async def process(self) -> str:
        """Process file and return extracted text."""
        processors = {
            'pdf': self._process_pdf,
            'docx': self._process_docx,
            'pptx': self._process_pptx,
            'txt': self._process_txt
        }
        
        processor = processors.get(self.file_type)
        if not processor:
            raise FileProcessingError(f"סוג קובץ לא נתמך: {self.file_type}")
            
        return await processor()
    
    async def _process_pdf(self) -> str:
        """Process PDF file."""
        try:
            return extract_text_from_pdf(self.file_path)
        except Exception as e:
            raise FileProcessingError(f"שגיאה בעיבוד קובץ PDF: {str(e)}")
    
    async def _process_docx(self) -> str:
        """Process DOCX file."""
        try:
            return extract_text_from_docx(self.file_path)
        except Exception as e:
            raise FileProcessingError(f"שגיאה בעיבוד קובץ DOCX: {str(e)}")
    
    async def _process_pptx(self) -> str:
        """Process PPTX file."""
        try:
            return extract_text_from_pptx(self.file_path)
        except Exception as e:
            raise FileProcessingError(f"שגיאה בעיבוד קובץ PPTX: {str(e)}")
    
    async def _process_txt(self) -> str:
        """Process TXT file."""
        try:
            return extract_text_from_txt(self.file_path)
        except Exception as e:
            raise FileProcessingError(f"שגיאה בעיבוד קובץ TXT: {str(e)}")

class TempFileManager:
    """Manages temporary files."""
    
    def __init__(self, temp_dir: str):
        self.temp_dir = temp_dir
        
    @contextmanager
    def temp_file(self, suffix: str) -> Generator[str, None, None]:
        """Create a temporary file that's automatically cleaned up."""
        temp_path = os.path.join(
            self.temp_dir,
            f"temp_{int(time.time())}_{random.randint(1000, 9999)}{suffix}"
        )
        try:
            yield temp_path
        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception as e:
                    logger.error(f"Failed to remove temp file {temp_path}: {e}")

# Initialize temp file manager
temp_manager = TempFileManager(TEMP_DIR)

# ====== Rate Limit Helpers ======
def check_rate_limit(user_id):
    """Check if user has exceeded rate limit."""
    current_time = time.time()
    user_times = user_request_times[user_id]

    # Remove old requests
    user_times = [t for t in user_times if current_time - t < RATE_LIMIT_PERIOD]
    user_request_times[user_id] = user_times

    if len(user_times) >= MAX_REQUESTS_PER_PERIOD:
        return False

    user_times.append(current_time)
    return True

# ====== Group Membership Check ======
async def user_can_use_service(bot: Bot, user_id: int) -> bool:
    """
    Checks if user can use the service based on daily usage
    """
    uses_today = await get_user_daily_uses(user_id)
    is_premium = await is_user_premium(user_id)
    daily_limit = PREMIUM_USES_PER_DAY if is_premium else FREE_USES_PER_DAY
    
    return uses_today < daily_limit

# ====== Command Handlers ======
  # Use this instead of the direct @retry decorator
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Start command handler with command overview."""
    user_id = update.effective_user.id
    
    # Check if this is a deep link with a quiz ID
    if context.args and context.args[0].startswith('quiz_'):
        try:
            quiz_id = int(context.args[0].split('_')[1])
            # Load and start the shared quiz
            with get_db_connection() as conn:
                c = conn.cursor()
                c.execute('''
                    SELECT title, questions_json 
                    FROM saved_quizzes 
                    WHERE quiz_id = ?
                ''', (quiz_id,))
                result = c.fetchone()
                
                if result:
                    title, questions_json = result
                    # Update play count
                    c.execute('''
                        UPDATE saved_quizzes 
                        SET times_played = times_played + 1 
                        WHERE quiz_id = ?
                    ''', (quiz_id,))
                    conn.commit()
                    
                    # Load questions and randomize options
                    questions = json.loads(questions_json)
                    randomized_questions = []
                    for question in questions:
                        randomized_questions.append(randomize_question_options(question))
                    
                    # Set up quiz in context
                    context.user_data["randomized_questions"] = randomized_questions
                    context.user_data["current_question_index"] = 0
                    context.user_data["correct_answers"] = 0
                    context.user_data[STATE_QUIZ_ACTIVE] = True
                    context.user_data['current_quiz_id'] = quiz_id
                    
                    await update.message.reply_text(
                        f"📝 *התחלת מבחן משותף*\n"
                        f"כותרת: {title}\n"
                        f"מתחילים\\! בהצלחה\\!",
                        parse_mode='MarkdownV2'
                    )
                    
                    # Start the quiz
                    await send_question(update, context)
                    return
                else:
                    await update.message.reply_text("❌ המבחן המבוקש לא נמצא.")
                    
        except (IndexError, ValueError) as e:
            logger.error(f"Error starting shared quiz: {e}")
            await update.message.reply_text("❌ קישור לא תקין.")
    
    # If not a shared quiz, show regular welcome message
    uses_today = await get_user_daily_uses(user_id)
    is_premium = await is_user_premium(user_id)
    daily_limit = PREMIUM_USES_PER_DAY if is_premium else FREE_USES_PER_DAY
    
    welcome_text = (
        f"ברוכים הבאים לסיכום.AI! 📚\n\n"
        f"נשארו לך {daily_limit - uses_today} ניסיונות להיום.\n\n"
        "🔍 איך זה עובד?\n"
        "שלחו לי את חומרי הלימוד שלכם ואני אהפוך אותם למבחן אמריקאי חכם!\n\n"
        "💡 פורמטים נתמכים:\n"
        "• PDF - מסמכים וספרים\n"
        "• DOCX/DOC - מסמכי וורד\n"
        "• TXT - קבצי טקסט\n"
        "• PPT/PPTX - מצגות\n\n"
        "📝 טיפים:\n"
        "• קבצים עד 50MB\n"
        "•  הטקסט חייב להיות ניתן לבחירה (לא כתב יד!)\n"
        "• ניתן לשמור מבחנים לשימוש חוזר\n\n"
        "🤖 פקודות נוספות:\n"
        "/help - מדריך מפורט\n"
        "/subscribe - שדרוג לפרימיום\n"
        "/list - מבחנים שמורים\n"
        "/support - תמיכה\n\n"
        "מוכנים להתחיל? שלחו קובץ! 🚀"
    )

    await update.message.reply_text(welcome_text)

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    user_id = update.effective_user.id
    uses_today = await get_user_daily_uses(user_id)
    is_premium = await is_user_premium(user_id)
    daily_limit = PREMIUM_USES_PER_DAY if is_premium else FREE_USES_PER_DAY

    help_text = (
        "🤖 *סיכום\\.AI \\- מדריך למשתמש*\n\n"
        "*פקודות בסיסיות:*\n"
        "• `/start` \\- התחלת שימוש בבוט\n"
        "• `/help` \\- הצגת מדריך זה\n"
        "• `/subscribe` \\- שדרוג לפרימיום\n"
        "• `/support` \\- יצירת קשר עם התמיכה\n\n"
        "*יצירת מבחנים:*\n"
        "• שלחו קובץ \\(PDF, DOCX, PPT וכו'\\)\n"
        "• המתינו ליצירת המבחן\n"
        "• ענו על השאלות וקבלו משוב מיידי\n"
        "• בסוף המבחן לחצו על 🔄 כדי להתחיל מחדש\n\n"
        "*ניהול מבחנים:*\n"
        "• `/save` \\- שמירת מבחן נוכחי\n"
        "• `/list` \\- הצגת מבחנים שמורים\n"
        "• `/play` \\- התחלת מבחן שמור\n"
        "• `/share` \\- לשתף את המבחן עם אחרים\n\n"
        "*סטטוס שימוש:*\n"
        f"• ניסיונות היום: {uses_today}/{daily_limit}\n"
        f"• סטטוס: {'פרימיום' if is_premium else 'רגיל'}\n\n"
        "*מגבלות:*\n"
        f"• עד {MAX_FILE_SIZE_MB}MB לקובץ\n"
        f"• {FREE_USES_PER_DAY} ניסיונות ביום למשתמשים רגילים\n"
        f"• {PREMIUM_USES_PER_DAY} ניסיונות ביום למשתמשי פרימיום\n"
    )

    await update.message.reply_text(help_text, parse_mode='MarkdownV2')

async def stats_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Show bot statistics to admins."""
    user_id = update.effective_user.id

    if user_id not in ADMIN_IDS:
        await update.message.reply_text("⛔️ פקודה זו זמינה רק למנהלי הבוט.")
        return

    stats_message = await get_stats_message()
    await update.message.reply_text(stats_message, parse_mode="Markdown")

async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle file uploads (PDF, DOCX, PPTX, etc)."""
    user_id = update.effective_user.id
    
    try:
        # Log the attempt
        logger.info(f"Document upload attempt by user {user_id}")
        
        # Check and increment usage
        uses_today = await get_user_daily_uses(user_id)
        is_premium = await is_user_premium(user_id)
        daily_limit = PREMIUM_USES_PER_DAY if is_premium else FREE_USES_PER_DAY
        
        if uses_today >= daily_limit:
            await update.message.reply_text(
                "השתמשת בכל הניסיונות שלך להיום.\n\n"
                "השתמש בפקודה /subscribe כדי לרכוש מנוי פרימיום ולקבל 30 ניסיונות ביום!"
            )
            return

        # Increment usage count first
        await increment_user_daily_uses(user_id)
        
        # Log the new count
        new_uses = await get_user_daily_uses(user_id)
        logger.info(f"User {user_id} daily uses increased from {uses_today} to {new_uses}")
        
        # Process file
        with temp_manager.temp_file(os.path.splitext(update.message.document.file_name)[1]) as temp_file:
            # Download file
            status_message = await update.message.reply_text("📥 מוריד את הקובץ...")
            if not await process_file(update, context, temp_file):
                await status_message.edit_text("❌ הורדת הקובץ נכשלה.")
                return

            # Process file
            await status_message.edit_text("📖 מעבד את הקובץ...")
            file_type = ALLOWED_EXTENSIONS[os.path.splitext(update.message.document.file_name.lower())[1]]
            
            processor = FileProcessor(temp_file, file_type)
            try:
                all_text = await processor.process()
            except FileProcessingError as e:
                await status_message.edit_text(f"❌ {str(e)}")
                return

            if not all_text.strip():
                await status_message.edit_text(
                    "❌ לא הצלחתי לחלץ טקסט קריא מהקובץ. "
                    "וודא שהטקסט במסמך ניתן לבחירה (לא רק תמונות)."
                )
                return

            # Generate quiz
            await status_message.edit_text("🧠 מייצר שאלות למבחן...")
            all_questions = await generate_quiz_questions(all_text)

            if not all_questions:
                await status_message.edit_text(
                    "❌ לא הצלחתי לייצר שאלות מקובץ זה. "
                    "וודא שיש בו מספיק טקסט קריא."
                )
                return

            # Trim or keep up to 20
            final_questions = (
                random.sample(all_questions, QUESTIONS_PER_QUIZ)
                if len(all_questions) >= QUESTIONS_PER_QUIZ
                else all_questions
            )

            # Save in context and begin quiz
            context.user_data["randomized_questions"] = final_questions
            context.user_data["original_questions"] = final_questions.copy()  # Save a copy of the original questions
            context.user_data["current_question_index"] = 0
            context.user_data["correct_answers"] = 0
            context.user_data[STATE_QUIZ_ACTIVE] = True

            # Cleanup
            await status_message.delete()
            await send_question(update, context)

            logger.info(f"User {user_id} started quiz from file {update.message.document.file_name}")

    except FileProcessingError as e:
        await handle_error(update, context, str(e))
    except DatabaseError as e:
        await handle_error(update, context, "שגיאת מסד נתונים")
        logger.error(f"Database error: {e}")
    except Exception as e:
        await handle_error(update, context, "שגיאה לא צפויה")
        logger.error(f"Unexpected error: {e}", exc_info=True)

# Add these optimization functions

# Create a separate function for formatting the question text without caching
def format_question_text(question_index, total_questions, question_text, options):
    """Format question text without caching."""
    try:
        return (
            f"*שאלה {question_index + 1} מתוך {total_questions}*\n\n"
            f"\u200F{question_text}\n\n"
            f"*תשובות אפשריות:*\n"
            f"\u200Fא. {options[0].split('. ', 1)[-1] if isinstance(options[0], str) else options[0]}\n"
            f"\u200Fב. {options[1].split('. ', 1)[-1] if isinstance(options[1], str) else options[1]}\n"
            f"\u200Fג. {options[2].split('. ', 1)[-1] if isinstance(options[2], str) else options[2]}\n"
            f"\u200Fד. {options[3].split('. ', 1)[-1] if isinstance(options[3], str) else options[3]}\n"
        )
    except Exception as e:
        logger.error(f"Error in format_question_text: {e}")
        return (
            f"*שאלה {question_index + 1} מתוך {total_questions}*\n\n"
            f"\u200F{question_text}\n\n"
            f"*תשובות אפשריות:*\n"
            f"א. ב. ג. ד. (שגיאה בטעינת אפשרויות)"
        )

# Simple cache dictionary for question texts
_question_cache = {}

def get_cached_question_text(question_index, total_questions, question_text, options):
    """Cache formatted question text using a manual cache instead of lru_cache."""
    # Create a simple hashable key using integers and strings only
    try:
        cache_key = (
            question_index, 
            total_questions, 
            question_text, 
            # Use only the first 20 chars of each option to create a hash
            tuple(str(opt)[:20] for opt in options)
        )
        
        # Check if the result is already in cache
        if cache_key in _question_cache:
            return _question_cache[cache_key]
            
        # If not in cache, format the text and store it
        result = format_question_text(question_index, total_questions, question_text, options)
        _question_cache[cache_key] = result
        
        # Keep cache size reasonable
        if len(_question_cache) > 200:  # Allow more cache entries than the original 128
            # Simple approach: clear half the cache when it gets too big
            keys_to_remove = list(_question_cache.keys())[:100]
            for k in keys_to_remove:
                _question_cache.pop(k, None)
                
        return result
    except Exception as e:
        logger.error(f"Cache error in get_cached_question_text: {e}")
        # Fall back to direct formatting without caching
        return format_question_text(question_index, total_questions, question_text, options)

# Enhanced file processing with ThreadPoolExecutor
async def process_file(update: Update, context: ContextTypes.DEFAULT_TYPE, temp_file: str) -> bool:
    """Process uploaded file with proper error handling."""
    try:
        # Get the file object
        file = await context.bot.get_file(update.message.document.file_id)
        
        # Download file directly (don't use ThreadPoolExecutor for coroutines)
        await file.download_to_drive(custom_path=temp_file)
        
        # Verify file exists
        if not os.path.exists(temp_file):
            logger.error(f"File download failed: {temp_file} doesn't exist")
            return False
            
        logger.info(f"Successfully downloaded file to {temp_file}, size: {os.path.getsize(temp_file)} bytes")
        return True
    except Exception as e:
        logger.error(f"File processing error: {e}", exc_info=True)
        return False

def extract_text_from_pdf(file_path):
    """Extract all text from a PDF file."""
    text_pages = []
    try:
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text and page_text.strip():
                        text_pages.append(page_text)
                except Exception as e:
                    logger.error(f"Error extracting text from page {page_num}: {str(e)}")
                    continue
    except Exception as e:
        logger.error(f"Error reading PDF file: {str(e)}")
    return "\n\n".join(text_pages)

def extract_text_from_docx(file_path):
    """Extract all text from a DOCX file using python-docx."""
    text_paragraphs = []
    try:
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                text_paragraphs.append(para.text.strip())
    except Exception as e:
        logger.error(f"Error reading DOCX file: {str(e)}")
    return "\n\n".join(text_paragraphs)

def extract_text_from_pptx(file_path):
    """Extract all text from a PowerPoint file."""
    text_parts = []
    try:
        prs = Presentation(file_path)
        
        # Extract title and text from each slide
        for slide in prs.slides:
            slide_text = []
            
            # Get text from shapes (including title)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Add slide text if not empty
            if slide_text:
                text_parts.append("\n".join(slide_text))
                
    except Exception as e:
        logger.error(f"Error reading PPTX file: {str(e)}")
        
    return "\n\n".join(text_parts)

def extract_text_from_txt(file_path):
    """Extract all text from a plain text file."""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
            return text.strip()
    except Exception as e:
        logger.error(f"Error reading TXT file: {str(e)}")
        return ""

# ====== Quiz Generation ======
def split_text_into_semantic_chunks(text: str, max_chunk_size: int = 2500) -> List[str]:
    """Split text into semantic chunks based on paragraphs."""
    paragraphs = text.split('\n\n')
    chunks = []
    current_chunk = []
    current_size = 0
    
    for paragraph in paragraphs:
        paragraph = paragraph.strip()
        if not paragraph:
            continue
            
        if current_size + len(paragraph) > max_chunk_size and current_chunk:
            chunks.append('\n\n'.join(current_chunk))
            current_chunk = []
            current_size = 0
            
        current_chunk.append(paragraph)
        current_size += len(paragraph)
    
    # Add the last chunk if it exists
    if current_chunk:
        chunks.append('\n\n'.join(current_chunk))
    
    # If no chunks were created, use the entire text as one chunk
    if not chunks and text.strip():
        chunks = [text.strip()]
    
    return chunks

# Add this new function after split_text_into_semantic_chunks function
def randomize_question_options(question):
    """
    Randomize the order of options in a question and update the correct_option_index.
    
    Args:
        question (dict): A question dictionary with 'options' and 'correct_option_index'
        
    Returns:
        dict: The question with randomized options and updated correct_option_index
    """
    if not isinstance(question, dict) or 'options' not in question or 'correct_option_index' not in question:
        return question
        
    try:
        # Get the current correct option
        original_options = list(question['options'])
        correct_option_index = question.get('correct_option_index', 0)
        
        if not (0 <= correct_option_index < len(original_options)):
            return question
            
        correct_option = original_options[correct_option_index]
        
        # Create pairs of (option, is_correct)
        option_pairs = [(opt, i == correct_option_index) for i, opt in enumerate(original_options)]
        
        # Shuffle the pairs
        random.shuffle(option_pairs)
        
        # Update the question with shuffled options
        question['options'] = [pair[0] for pair in option_pairs]
        
        # Find the new index of the correct option
        for i, (_, is_correct) in enumerate(option_pairs):
            if is_correct:
                question['correct_option_index'] = i
                break
                
        return question
    except Exception as e:
        logger.error(f"Error randomizing options: {e}")
        return question

# More efficient quiz generation
async def generate_quiz_questions(text: str, max_retries: int = 3) -> List[Dict]:
    """Generate quiz questions with better parallelization and caching."""
    chunks = split_text_into_semantic_chunks(text, max_chunk_size=2500)
    all_questions = []
    questions_per_chunk = math.ceil(20 / max(len(chunks), 1))
    
    # Increase concurrent API requests from 2 to 4 to utilize both CPU cores
    semaphore = asyncio.Semaphore(4)
    
    # Add timeout to API calls to prevent hanging
    async def process_chunk(chunk):
        async with semaphore:
            for attempt in range(max_retries):
                try:
                    # Use asyncio.wait_for to add a timeout
                    prompt = (
                        f"Create exactly {questions_per_chunk} multiple choice questions in Hebrew that assess mastery "
                        "of the concepts from the background content. Generate questions that could be answered by someone "
                        "who truly understands the material, without needing to reference specific text.\n\n"
                        "CRITICAL RULES:\n"
                        "1. NEVER use phrases like 'according to the text', 'based on the passage', or any direct text references\n"
                        "2. Questions must be in Hebrew\n"
                        "3. Each question must have exactly 4 options labeled as: א, ב, ג, ד\n"
                        "4. The correct answer must be unambiguously correct and fully supported by the background content\n"
                        "5. Focus on testing:\n"
                        "   - Deep comprehension of concepts\n"
                        "   - Ability to apply principles\n"
                        "   - Understanding of relationships and implications\n"
                        "   - Critical thinking about the subject matter\n"
                        "7. Each explanation must clearly justify why the correct answer is the only valid choice\n"
                        "8. Do not use trailing commas in arrays\n\n"
                        "VERIFICATION STEPS:\n"
                        "- Verify each question can be answered without seeing the original text\n"
                        "- Confirm no direct references to the source material\n"
                        "- Ensure correct answers are definitively supported by the background content\n\n"                        
                        "Return a valid JSON array where each question has this exact format:\n"
                        '{\n'
                        '    "question": "שאלה בעברית?",\n'
                        '    "options": ["א. אפשרות 1", "ב. אפשרות 2", "ג. אפשרות 3", "ד. אפשרות 4"],\n'
                        '    "correct_option_index": 0,\n'
                        '    "explanation": "הסבר קצר"\n'
                        '}\n\n'
                        f"Background content to derive concepts from:\n{chunk}"
                    )
                    
                    # Add a timeout to the API call to prevent hanging
                    response = model.generate_content(prompt)
                    
                    if not response or not response.text:
                        continue
                    
                    # Clean the response text
                    cleaned_text = response.text.strip()
                    if '```json' in cleaned_text:
                        cleaned_text = cleaned_text.split('```json')[1].split('```')[0].strip()
                    elif '```' in cleaned_text:
                        cleaned_text = cleaned_text.split('```')[1].split('```')[0].strip()
                    
                    # Enhanced JSON repair for common AI-generated JSON errors
                    import re
                    # Fix trailing commas in arrays (including nested arrays)
                    cleaned_text = re.sub(r',\s*\]', ']', cleaned_text)
                    # Fix trailing commas in objects (including nested objects)
                    cleaned_text = re.sub(r',\s*\}', '}', cleaned_text)
                    # Fix missing commas between array items
                    cleaned_text = re.sub(r'"\s*"', '", "', cleaned_text)
                    # Fix missing closing brackets for options array
                    cleaned_text = re.sub(r'(".+?"),?\s*"correct_option_index"', r'\1], "correct_option_index"', cleaned_text)
                    # Add missing closing brackets for nested arrays
                    bracket_count = cleaned_text.count('[') - cleaned_text.count(']')
                    if bracket_count > 0:
                        cleaned_text += ']' * bracket_count
                    
                    # Try advanced JSON repair if standard fixes don't work
                    try:
                        chunk_questions = json.loads(cleaned_text)
                        if not isinstance(chunk_questions, list):
                            chunk_questions = [chunk_questions]
                            
                        # Randomize options for each question
                        for i in range(len(chunk_questions)):
                            chunk_questions[i] = randomize_question_options(chunk_questions[i])
                            
                        return chunk_questions
                    except json.JSONDecodeError as e:
                        logger.error(f"JSON parsing error: {e}")
                        logger.error(f"Cleaned text: {cleaned_text}")
                        
                        # Last resort - try repair_json function
                        try:
                            repaired_json = repair_json(cleaned_text)
                            chunk_questions = json.loads(repaired_json)
                            if not isinstance(chunk_questions, list):
                                chunk_questions = [chunk_questions]
                                
                            # Randomize options for each question
                            for i in range(len(chunk_questions)):
                                chunk_questions[i] = randomize_question_options(chunk_questions[i])
                                
                            logger.info("Successfully repaired malformed JSON")
                            return chunk_questions
                        except Exception as repair_error:
                            logger.error(f"JSON repair failed: {repair_error}")
                            continue
                except Exception as e:
                    logger.error(f"Error in chunk {chunks.index(chunk)}, attempt {attempt+1}: {e}")
                    if "429" in str(e):
                        await asyncio.sleep(5 * (attempt + 1))
                    continue
            return []
    
    # Process chunks in parallel
    tasks = [process_chunk(chunk) for chunk in chunks]
    results = await asyncio.gather(*tasks)
    
    # Combine all questions
    for questions in results:
        all_questions.extend(questions)
    
    # Return results
    if len(all_questions) >= 20:
        return random.sample(all_questions, 20)
    return all_questions

# ====== Quiz Flow Handlers ======
async def send_question(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Sends the current question with optimized performance."""
    try:
        questions = context.user_data.get("randomized_questions", [])
        current_q = context.user_data.get("current_question_index", 0)
        
        # Check if we've reached the end of the quiz
        if not questions or current_q >= len(questions):
            await show_result(update, context)
            return
            
        question = questions[current_q]
        q_index = current_q
        
        # Debug info for question structure
        logger.debug(f"Question structure: {type(question)}")
        logger.debug(f"Question keys: {question.keys() if isinstance(question, dict) else 'Not a dict'}")
        
        if 'options' in question:
            logger.debug(f"Options type: {type(question['options'])}")
            logger.debug(f"Options value: {str(question['options'])[:100]}")
            
            # Pre-process options if needed - ensure it's a proper list of strings
            # This is a safeguard against malformed data
            if isinstance(question['options'], list):
                processed_options = []
                for opt in question['options']:
                    if isinstance(opt, (list, dict)):
                        # Convert complex structures to string
                        processed_options.append(str(opt))
                    else:
                        processed_options.append(opt)
                question['options'] = processed_options
        
        # Create keyboard with minimal processing
        keyboard = [
            [InlineKeyboardButton(text=letter, callback_data=f"{q_index}|{idx}")]
            for idx, letter in enumerate(HEBREW_LETTERS[:4])
        ]
        
        # Use the cached question text formatter for better performance
        try:
            question_text = get_cached_question_text(
                current_q, 
                len(questions), 
                question['question'], 
                question['options']
            )
        except Exception as e:
            logger.error(f"Error formatting question text: {e}")
            # Fallback to simple formatting
            question_text = f"*שאלה {current_q + 1} מתוך {len(questions)}*\n\n{question['question']}"
        
        # Simplified response handling
        if update.callback_query:
            try:
                # Use a reasonable timeout for edit operations
                await update.callback_query.edit_message_text(
                    text=question_text,
                    reply_markup=InlineKeyboardMarkup(keyboard),
                    parse_mode='Markdown'
                )
            except BadRequest as e:
                # Only critical errors should be re-raised
                if "Message is not modified" not in str(e) and "Message to edit not found" not in str(e):
                    logger.error(f"Error updating message: {e}")
                    # Send a new message instead of failing
                    await context.bot.send_message(
                        chat_id=update.effective_chat.id,
                        text=question_text,
                        reply_markup=InlineKeyboardMarkup(keyboard),
                        parse_mode='Markdown'
                    )
        else:
            await update.message.reply_text(
                text=question_text,
                reply_markup=InlineKeyboardMarkup(keyboard),
                parse_mode='Markdown'
            )
            
    except Exception as e:
        logger.error(f"Error in send_question: {e}")
        # Always try to recover by sending a new message
        try:
            if update.effective_chat:
                await context.bot.send_message(
                    chat_id=update.effective_chat.id,
                    text="⚠️ אירעה שגיאה בהצגת השאלה. מנסה להמשיך...",
                    reply_markup=InlineKeyboardMarkup([
                        [InlineKeyboardButton("המשך", callback_data=f"continue")]
                    ])
                )
        except:
            pass

@with_network_retry
async def handle_answer(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle button clicks from users (answer selections)."""
    query = update.callback_query
    
    try:
        # Safety check
        if not query:
            return
            
        await query.answer()  # Acknowledge the button click to remove the loading indicator
        
        # Handle "next" button to go to the next question
        if query.data == "next":
            await send_question(update, context)
            return
        
        # Handle "redo_quiz" button to restart the quiz with same questions but randomized
        if query.data == "redo_quiz":
            # Get original questions from context
            original_questions = context.user_data.get("original_questions", [])
            
            if not original_questions:
                # Fallback if original questions not found
                await query.edit_message_text(
                    "⚠️ לא ניתן להתחיל מחדש את המבחן. נסה להריץ מבחן חדש.",
                    parse_mode='Markdown'
                )
                logger.error(f"Redo quiz failed: No original questions found for user {update.effective_user.id}")
                return
                
            # Randomize options for each question
            randomized_questions = []
            for question in original_questions:
                randomized_questions.append(randomize_question_options(question))
            
            # Log the restart event
            logger.info(f"User {update.effective_user.id} restarting quiz with {len(randomized_questions)} questions")
                
            # Reset quiz state
            context.user_data["randomized_questions"] = randomized_questions
            context.user_data["current_question_index"] = 0
            context.user_data["correct_answers"] = 0
            context.user_data[STATE_QUIZ_ACTIVE] = True
            context.user_data.pop("quiz_completed", None)
            
            # Start the quiz again
            await query.edit_message_text(
                "🔄 *המבחן מתחיל מחדש\\!*\n\nהשאלות זהות אך סדר התשובות שונה\\.",
                parse_mode='MarkdownV2'
            )
            await send_question(update, context)
            return
        
        # Handle "save_quiz" button to save the quiz
        if query.data == "save_quiz":
            await query.edit_message_text(
                "אנא שלח שם למבחן (או /cancel לביטול)",
                parse_mode='Markdown'
            )
            context.user_data['waiting_for_quiz_title'] = True
            context.user_data['quiz_to_save'] = context.user_data.get("randomized_questions", [])
            return
        
        # Handle "list_quizzes" button to show saved quizzes
        if query.data == "list_quizzes":
            # Create a new Update object with the user's message for list_saved_quizzes
            new_update = Update(update.update_id, message=update.effective_message)
            await list_saved_quizzes(new_update, context)
            return
        
        # Handle "share_quiz" button
        if query.data == "share_quiz":
            # Check if current quiz is already saved
            quiz_id = context.user_data.get('current_quiz_id')
            if quiz_id:
                # If quiz is already saved, share it directly
                await share_quiz(update, context, quiz_id)
            else:
                # If not saved, prompt to save first
                await query.edit_message_text(
                    "עליך לשמור את המבחן תחילה. אנא שלח שם למבחן (או /cancel לביטול)",
                    parse_mode='Markdown'
                )
                context.user_data['waiting_for_quiz_title'] = True
                context.user_data['share_after_save'] = True
                context.user_data['quiz_to_save'] = context.user_data.get("randomized_questions", [])
            return
        
        # FIXED: Handle special callback data formats
        if query.data == "continue":
            # This is a recovery button - just show the current question
            await send_question(update, context)
            return
            
        # Handle regular answer selection
        try:
            q_index, selected_option = map(int, query.data.split('|'))
            
            # Get minimal required data
            questions = context.user_data.get("randomized_questions", [])
            if not questions or q_index >= len(questions):
                return
                
            current_question = questions[q_index]
            correct_option = current_question.get("correct_option_index", 0)
            
            # Update score first (fast operation)
            if selected_option == correct_option:
                context.user_data["correct_answers"] = context.user_data.get("correct_answers", 0) + 1
                is_correct = True
            else:
                is_correct = False
            
            # Move to next question (but don't show it yet - wait for button press)
            context.user_data["current_question_index"] = q_index + 1
            
            # Show feedback in the same message - modify the current message
            try:
                # Get feedback text
                correct_answer = current_question['options'][correct_option].split('. ', 1)[-1]
                
                # Include the original question in the feedback
                question_text = current_question['question']
                options_text = '\n'.join(current_question['options'])
                
                # Format feedback with question included
                feedback = f"*שאלה {q_index + 1}*\n\n{question_text}\n\n{options_text}\n\n"
                
                if is_correct:
                    feedback += f"✅ *תשובה נכונה!*\nהתשובה הנכונה: {correct_answer}"
                else:
                    selected_text = current_question['options'][selected_option].split('. ', 1)[-1]
                    feedback += f"❌ *טעות!*\nבחרת: {selected_text}\nהתשובה הנכונה: {correct_answer}"
                
                # Add explanation if available
                if 'explanation' in current_question and current_question['explanation']:
                    feedback += f"\n\n*הסבר:* {current_question['explanation']}"
                
                # Create next button
                next_keyboard = InlineKeyboardMarkup([
                    [InlineKeyboardButton("→ לשאלה הבאה", callback_data="next")]
                ])
                
                # Edit message with feedback
                await query.edit_message_text(
                    text=feedback,
                    reply_markup=next_keyboard,
                    parse_mode='Markdown'
                )
                
            except Exception as e:
                logger.error(f"Error showing feedback: {e}")
                # Continue to next question if feedback fails
                await send_question(update, context)
                
        except (ValueError, IndexError) as e:
            logger.error(f"Invalid callback data: {query.data}, Error: {e}")
            
        except Exception as e:
            logger.error(f"Error processing answer: {e}")
            await query.edit_message_text(
                "⚠️ אירעה שגיאה בעיבוד התשובה. נסה להמשיך.",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("המשך", callback_data="continue")]
                ])
            )
        
    except Exception as e:
        logger.error(f"Error handling answer: {e}", exc_info=True)

async def send_feedback_message(bot, user_id, question, selected_option, with_next_button=False):
    """Send feedback message with optional next button."""
    try:
        correct_option = question.get("correct_option_index", 0)
        correct_answer = question['options'][correct_option].split('. ', 1)[-1]
        
        if selected_option == correct_option:
            response = f"✅ תשובה נכונה!\nהתשובה הנכונה: {correct_answer}"
        else:
            selected_text = question['options'][selected_option].split('. ', 1)[-1]
            response = f"❌ טעות! בחרת: {selected_text}\nהתשובה הנכונה: {correct_answer}"
        
        # Add explanation if available
        if 'explanation' in question and question['explanation']:
            response += f"\n\nהסבר: {question['explanation']}"
            
        # Add next button if requested
        if with_next_button:
            keyboard = InlineKeyboardMarkup([
                [InlineKeyboardButton("→ לשאלה הבאה", callback_data="next")]
            ])
            
            # Send message with next button (don't delete it)
            message = await bot.send_message(
                chat_id=user_id,
                text=response,
                reply_markup=keyboard,
                parse_mode='Markdown',
                disable_notification=True
            )
            return message
        else:
            # Send temporary message that will be deleted after a short time
            message = await bot.send_message(
                chat_id=user_id,
                text=response + "\n\n(הודעה זו תימחק אוטומטית)",
                disable_notification=True
            )
            
            # Schedule deletion after 3 seconds
            async def delete_message():
                await asyncio.sleep(3)
                try:
                    await message.delete()
                except Exception as e:
                    logger.debug(f"Couldn't delete message: {e}")
                    
            # Run deletion in background
            asyncio.create_task(delete_message())
            
    except Exception as e:
        logger.error(f"Error sending feedback: {e}")
        return None

async def show_result(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Display final quiz result in an improved UX format."""
    try:
        user_id = update.effective_user.id
        correct_answers = context.user_data.get("correct_answers", 0)
        questions = context.user_data.get("randomized_questions", [])
        total = len(questions)
        score = (correct_answers / total) * 100 if total > 0 else 0
        
        # Format message with proper escaping
        def escape_markdown_v2(text):
            """Escape special characters for MarkdownV2."""
            special_chars = '_*[]()~`>#+-=|{}.!'
            return ''.join(f'\\{c}' if c in special_chars else c for c in str(text))
        
        formatted_score = escape_markdown_v2(f"{score:.1f}")
        quiz_id = context.user_data.get('current_quiz_id', '')
        
        # Add emoji based on score
        score_emoji = "🎉" if score >= 90 else "🌟" if score >= 70 else "👍" if score >= 50 else "😐"
        
        score_text = (
            f"{score_emoji} *התוצאה שלך* {score_emoji}\n\n"
            f"ענית נכון על {escape_markdown_v2(correct_answers)} מתוך {escape_markdown_v2(total)} שאלות\n"
            f"ציון סופי: {formatted_score}\\%\n\n"
            f"*מה תרצה לעשות עכשיו?*"
        )
        
        # Create keyboard with all buttons - redo, save, list, share
        redo_keyboard = InlineKeyboardMarkup([
            [InlineKeyboardButton("🔄 התחל את המבחן מחדש", callback_data="redo_quiz")],
            [InlineKeyboardButton("💾 שמור את המבחן", callback_data="save_quiz")],
            [InlineKeyboardButton("📋 מבחנים שמורים", callback_data="list_quizzes")],
            [InlineKeyboardButton("🔗 שתף את המבחן", callback_data="share_quiz")]
        ])
        
        # Try to edit the existing message first if there's a callback query
        if update.callback_query:
            try:
                await update.callback_query.edit_message_text(
                    text=score_text,
                    reply_markup=redo_keyboard,
                    parse_mode='MarkdownV2'
                )
                return
            except Exception as e:
                logger.debug(f"Could not edit message for result, sending new one: {e}")
                
        # If editing fails or no callback query exists (e.g., when directly called), send a new message
        await context.bot.send_message(
            chat_id=user_id,
            text=score_text,
            reply_markup=redo_keyboard,
            parse_mode='MarkdownV2'
        )
        
        # but keep the quiz data for saving/sharing
        context.user_data["quiz_completed"] = True
        context.user_data["final_score"] = score
        # Store original questions for potential redo
        context.user_data["original_questions"] = questions.copy()
        
        update_user_stats(user_id, quiz_completed=True, score=int(score), total_questions=total)
        
    except Exception as e:
        logger.error(f"Error in show_result: {str(e)}")
        await handle_error(update, context, "הצגת התוצאות נכשלה")

# ====== Error Handler ======
async def error_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Enhanced error handler with timeouts awareness."""
    try:
        # Handle callback timeouts specifically
        if isinstance(context.error, BadRequest) and any(phrase in str(context.error) for phrase in 
                                                        ["Query is too old", "Message to edit not found", "Message can't be edited"]):
            if update and update.callback_query:
                user_id = update.callback_query.from_user.id
                chat_id = update.effective_chat.id
                
                # Send new message instead of editing
                try:
                    await context.bot.send_message(
                        chat_id=chat_id,
                        text="⏱️ הפעולה נמשכה זמן רב מדי. המבחן ממשיך...",
                        reply_markup=InlineKeyboardMarkup([
                            [InlineKeyboardButton("המשך במבחן", callback_data="continue")]
                        ])
                    )
                    return
                except Exception as inner_e:
                    logger.error(f"Failed in error recovery: {inner_e}")
        
        # Rest of your existing error handler...
    except Exception as e:
        logger.error(f"Error in error handler: {e}", exc_info=True)

# ====== Statistics Functions ======
def update_user_stats(user_id: int, quiz_completed: bool = False, score: int = 0, total_questions: int = 0) -> None:
    """Update user statistics in SQLite database."""
    with get_db_connection() as conn:
        c = conn.cursor()

        try:
            # Convert date to string in ISO format
            today_str = date.today().isoformat()

            # Update users table
            c.execute('''
                INSERT OR IGNORE INTO users (user_id, first_seen)
                VALUES (?, ?)
            ''', (user_id, today_str))

            # Update daily usage
            c.execute('''
                INSERT OR IGNORE INTO daily_usage (date, user_id)
                VALUES (?, ?)
            ''', (today_str, user_id))

            # If quiz completed, add to quizzes table
            if quiz_completed:
                # Convert datetime to string
                now_str = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                c.execute('''
                    INSERT INTO quizzes (user_id, timestamp, score, total_questions)
                    VALUES (?, ?, ?, ?)
                ''', (user_id, now_str, score, total_questions))

            conn.commit()

        except Exception as e:
            logger.error(f"Database error in update_user_stats: {e}")

async def get_stats_message() -> str:
    """Get formatted statistics message."""
    with get_db_connection() as conn:
        c = conn.cursor()
        
        try:
            # Get total unique users from daily_usage (more accurate than users table)
            c.execute('SELECT COUNT(DISTINCT user_id) FROM daily_usage')
            total_users = c.fetchone()[0]

            # Get today's active users
            today = date.today().isoformat()
            c.execute('''
                SELECT COUNT(DISTINCT user_id) 
                FROM daily_usage 
                WHERE date = ?
            ''', (today,))
            active_today = c.fetchone()[0]

            # Get total tests (sum of all attempts)
            c.execute('SELECT SUM(attempts) FROM daily_usage')
            total_tests = c.fetchone()[0] or 0

            # Get today's tests
            c.execute('''
                SELECT SUM(attempts) 
                FROM daily_usage 
                WHERE date = ?
            ''', (today,))
            tests_today = c.fetchone()[0] or 0

            # Get total premium users
            c.execute('''
                SELECT COUNT(*) FROM subscriptions 
                WHERE subscribed_until > date('now')
            ''')
            premium_users = c.fetchone()[0]

            stats_message = (
                "*📊 סטטיסטיקות הבוט*\n\n"
                f"👥 *סה״כ משתמשים:* {total_users}\n"
                f"👤 *משתמשים היום:* {active_today}\n"
                f"📚 *סה״כ מבחנים:* {total_tests}\n"
                f"📖 *מבחנים היום:* {tests_today}\n"
                f"⭐️ *משתמשי פרימיום:* {premium_users}"
            )
            
            return stats_message

        except Exception as e:
            logger.error(f"Database error in get_stats_message: {str(e)}")
            return "❌ שגיאה בטעינת הסטטיסטיקות"

def init_db():
    """Initialize database with all required tables and optimize for concurrent access."""
    with get_db_connection() as conn:
        c = conn.cursor()
        
        # Enable WAL mode and other optimizations at DB level
        c.execute("PRAGMA journal_mode = WAL")
        c.execute("PRAGMA synchronous = NORMAL")
        c.execute("PRAGMA temp_store = MEMORY")
        c.execute("PRAGMA cache_size = -2000")  # Use ~2MB of memory for cache
        c.execute("PRAGMA busy_timeout = 1000")  # Wait up to 1 second when db is busy
        c.execute("PRAGMA mmap_size = 30000000")  # Use memory-mapped I/O for better performance
        
        # Create users table
        c.execute('''
            CREATE TABLE IF NOT EXISTS users (
                user_id INTEGER PRIMARY KEY,
                first_seen DATE,
                username TEXT,
                last_active DATE
            )
        ''')
        
        # Create subscriptions table
        c.execute('''
            CREATE TABLE IF NOT EXISTS subscriptions (
                user_id INTEGER PRIMARY KEY,
                subscribed_until DATE,
                subscription_type TEXT,
                FOREIGN KEY (user_id) REFERENCES users(user_id)
            )
        ''')
        
        # Create saved_quizzes table (without dropping it)
        c.execute('''
            CREATE TABLE IF NOT EXISTS saved_quizzes (
                quiz_id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                title TEXT,
                questions_json TEXT,
                created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                times_played INTEGER DEFAULT 0,
                FOREIGN KEY (user_id) REFERENCES users(user_id)
            )
        ''')
        
        # Create daily_usage table
        c.execute('''
            CREATE TABLE IF NOT EXISTS daily_usage (
                user_id INTEGER,
                date DATE,
                attempts INTEGER DEFAULT 1,
                PRIMARY KEY (user_id, date),
                FOREIGN KEY (user_id) REFERENCES users(user_id)
            )
        ''')
        
        # Create quizzes table for tracking completed quizzes
        c.execute('''
            CREATE TABLE IF NOT EXISTS quizzes (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id INTEGER,
                score REAL,
                total_questions INTEGER,
                timestamp DATETIME DEFAULT CURRENT_TIMESTAMP,
                FOREIGN KEY (user_id) REFERENCES users(user_id)
            )
        ''')
        
        conn.commit()
        
        # Run optimize_db to create all necessary indexes
        optimize_db()

def optimize_db():
    """Create indexes and optimize database for performance."""
    with get_db_connection() as conn:
        c = conn.cursor()
        
        # Basic indexes for performance (existing)
        c.execute('CREATE INDEX IF NOT EXISTS idx_daily_usage_date ON daily_usage(date)')
        c.execute('CREATE INDEX IF NOT EXISTS idx_quizzes_user_id ON quizzes(user_id)')
        
        # Additional indexes for better performance under load
        c.execute('CREATE INDEX IF NOT EXISTS idx_saved_quizzes_user_id ON saved_quizzes(user_id)')
        c.execute('CREATE INDEX IF NOT EXISTS idx_saved_quizzes_created ON saved_quizzes(created_at)')
        c.execute('CREATE INDEX IF NOT EXISTS idx_subscriptions_user ON subscriptions(user_id)')
        c.execute('CREATE INDEX IF NOT EXISTS idx_subscriptions_expiry ON subscriptions(subscribed_until)')
        c.execute('CREATE INDEX IF NOT EXISTS idx_users_last_active ON users(last_active)')
        c.execute('CREATE INDEX IF NOT EXISTS idx_quizzes_timestamp ON quizzes(timestamp)')
        
        # Run ANALYZE to update statistics
        c.execute('ANALYZE')
        
        conn.commit()

def vacuum_db():
    """Perform VACUUM operation to optimize database size and performance."""
    # Connect directly to enable vacuum (can't be done with WAL)
    conn = sqlite3.connect(DB_PATH)
    try:
        # Temporarily disable WAL to perform VACUUM
        conn.execute("PRAGMA journal_mode = DELETE")
        conn.execute("VACUUM")
        conn.execute("PRAGMA journal_mode = WAL")  # Switch back to WAL
        conn.execute("PRAGMA optimize")  # New in SQLite 3.18.0
        conn.commit()
        logger.info("Database vacuum completed successfully")
    except Exception as e:
        logger.error(f"Error during database vacuum: {e}")
    finally:
        conn.close()

# Optimize database connections with connection pooling
import sqlite3
import threading
import asyncio
from functools import wraps

# Global connection dictionary for basic pooling
_connection_pool = {}
_pool_lock = threading.Lock()
_last_used = {}  # Add tracking of last usage
# Set maximum concurrent database connections
MAX_DB_CONNECTIONS = 15
_connection_semaphore = asyncio.Semaphore(MAX_DB_CONNECTIONS)

@contextmanager
async def get_async_db_connection():
    """Async version of get_db_connection with semaphore limiting."""
    try:
        # Acquire the semaphore to limit total concurrent database connections
        await _connection_semaphore.acquire()
        
        # Now use the synchronous connection manager which handles thread-based pooling
        with get_db_connection() as conn:
            yield conn
    finally:
        # Release the semaphore when done
        _connection_semaphore.release()

@contextmanager
def get_db_connection():
    """Improved connection pooling for database access with connection validation."""
    thread_id = threading.get_ident()
    
    with _pool_lock:
        current_time = time.time()
        
        # Check if connection exists and is valid
        if thread_id in _connection_pool:
            conn = _connection_pool[thread_id]
            try:
                # Test if connection is still valid with a simple query
                conn.execute("SELECT 1")
                _last_used[thread_id] = current_time
            except (sqlite3.ProgrammingError, sqlite3.OperationalError):
                # Connection is closed or invalid, create a new one
                logger.info(f"Replacing invalid connection for thread {thread_id}")
                try:
                    conn.close()
                except:
                    pass
                conn = None
        else:
            conn = None
            
        # Create new connection if needed
        if conn is None:
            conn = sqlite3.connect(
                DB_PATH, 
                timeout=5.0,
                isolation_level=None,  # Enable autocommit mode
                check_same_thread=False  # Allow cross-thread usage
            )
            # Enable WAL mode for better concurrency
            conn.execute('PRAGMA journal_mode = WAL')
            # Optimize performance
            conn.execute('PRAGMA synchronous = NORMAL')
            conn.execute('PRAGMA temp_store = MEMORY')
            conn.execute('PRAGMA cache_size = -2000')  # Use ~2MB of memory for cache
            conn.execute('PRAGMA busy_timeout = 1000')  # Wait up to 1 second when db is busy
            
            conn.row_factory = sqlite3.Row  # Enable row factory for better results
            
            _connection_pool[thread_id] = conn
            _last_used[thread_id] = current_time
            logger.debug(f"Created new DB connection for thread {thread_id}")
    
    try:
        yield conn
    except Exception as e:
        logger.error(f"Database error: {e}")
        # Don't close on error, just propagate the exception
        raise

# Add decorator for atomic transactions
def atomic_transaction(func):
    """Decorator to ensure database operations run in atomic transactions."""
    @wraps(func)
    async def wrapper(*args, **kwargs):
        async with get_async_db_connection() as conn:
            try:
                conn.execute("BEGIN IMMEDIATE")
                result = await func(conn, *args, **kwargs)
                conn.execute("COMMIT")
                return result
            except Exception as e:
                conn.execute("ROLLBACK")
                logger.error(f"Transaction error in {func.__name__}: {e}")
                raise
    return wrapper

# Add retry function for database operations
async def resilient_db_operation(operation, fallback=None, max_retries=3):
    """Execute database operations with retry logic for handling database locks."""
    for attempt in range(max_retries):
        try:
            return await operation()
        except sqlite3.OperationalError as e:
            if "database is locked" in str(e) and attempt < max_retries - 1:
                wait_time = 0.2 * (2 ** attempt)  # Exponential backoff
                logger.warning(f"Database locked, retrying in {wait_time:.2f}s (attempt {attempt+1}/{max_retries})")
                await asyncio.sleep(wait_time)
            else:
                logger.error(f"Database operation failed after {attempt+1} attempts: {e}")
                if fallback:
                    return fallback()
                raise

async def save_quiz(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Save the current quiz for future replay."""
    user_id = update.effective_user.id

    # Check if there's an active or completed quiz
    if not (context.user_data.get(STATE_QUIZ_ACTIVE) or context.user_data.get("quiz_completed")):
        await update.message.reply_text("אין מבחן פעיל לשמירה. צור קודם מבחן חדש!")
        return
    
    questions = context.user_data.get("randomized_questions", [])
    if not questions:
        await update.message.reply_text("לא נמצאו שאלות לשמירה!")
        return

    # If title was provided with command, save directly
    if context.args:
        title = " ".join(context.args)
        if validate_quiz_title(title):
            await save_quiz_with_title(update, context, title)
        else:
            await update.message.reply_text(
                "❌ שם המבחן אינו תקין. אנא בחר שם באורך של 1-100 תווים."
            )
    else:
        # Otherwise, ask for title interactively
        await update.message.reply_text("אנא שלח שם למבחן (או /cancel לביטול)")
        context.user_data['waiting_for_quiz_title'] = True
        context.user_data['quiz_to_save'] = questions

async def handle_quiz_title(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle the quiz title input with validation."""
    title = update.message.text
    
    if not validate_quiz_title(title):
        await update.message.reply_text(
            "❌ שם המבחן אינו תקין. אנא בחר שם באורך של 1-100 תווים."
        )
        return
    
    await save_quiz_with_title(update, context, title)

async def save_quiz_with_title(update: Update, context: ContextTypes.DEFAULT_TYPE, title: str) -> None:
    """Helper function to save quiz with given title."""
    questions = context.user_data.get("randomized_questions", [])
    share_after_save = context.user_data.get('share_after_save', False)
    
    with get_db_connection() as conn:
        c = conn.cursor()

        try:
            # Debug logging
            logger.info(f"Saving quiz with title: {title}")
            logger.info(f"Questions count: {len(questions)}")

            c.execute('''
                INSERT INTO saved_quizzes (user_id, title, questions_json, created_at)
                VALUES (?, ?, ?, ?)
            ''', (update.effective_user.id, title, json.dumps(questions), datetime.now().strftime('%Y-%m-%d %H:%M:%S')))

            quiz_id = c.lastrowid
            conn.commit()
            
            # Store the quiz_id in context for future use
            context.user_data['current_quiz_id'] = quiz_id

            success_message = (
                f"✅ המבחן נשמר בהצלחה!\n"
                f"כותרת: {title}\n"
                f"מזהה: {quiz_id}\n\n"
                f"להפעלת המבחן השתמש בפקודה:\n"
                f"/play {quiz_id}"
            )
            
            await update.message.reply_text(success_message)

            # Clear the waiting state
            context.user_data.pop('waiting_for_quiz_title', None)
            context.user_data.pop('quiz_to_save', None)
            
            # If share_after_save is True, share the quiz immediately after saving
            if share_after_save:
                context.user_data.pop('share_after_save', None)
                await share_quiz(update, context, quiz_id)

        except Exception as e:
            logger.error(f"Error saving quiz: {str(e)}")
            await update.message.reply_text("❌ שגיאה בשמירת המבחן. אנא נסה שוב.")

async def list_saved_quizzes(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """List all saved quizzes for the current user with caching for better performance."""
    user_id = update.effective_user.id
    
    try:
        # Use cached quiz list if available
        quizzes = await db_get_user_quizzes_with_caching(user_id)
        
        if not quizzes:
            await update.message.reply_text("לא נמצאו מבחנים שמורים!")
            return

        def escape_markdown(text):
            """Helper function to escape special characters for MarkdownV2"""
            special_chars = ['_', '*', '[', ']', '(', ')', '~', '`', '>', '#', '+', '-', '=', '|', '{', '}', '.', '!']
            for char in special_chars:
                text = text.replace(char, f"\\{char}")
            return text

        message = "📚 *המבחנים השמורים שלך*\n\n"
        for quiz in quizzes:
            safe_title = escape_markdown(str(quiz['title']))
            safe_created_at = escape_markdown(str(quiz['created_at']))

            message += (
                f"*מזהה מבחן:* `{quiz['quiz_id']}`\n"
                f"*כותרת:* {safe_title}\n"
                f"*נוצר בתאריך:* {safe_created_at}\n"
                f"*מספר פעמים ששוחק:* {quiz['times_played']}\n"
                f"השתמש בפקודה `/play {quiz['quiz_id']}` להפעלת המבחן\n\n"
            )

        await update.message.reply_text(message, parse_mode="MarkdownV2")

    except Exception as e:
        logger.error(f"Error listing quizzes: {str(e)}")
        await update.message.reply_text("❌ שגיאה בטעינת המבחנים. אנא נסה שוב.")

async def play_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """First step of playing a saved quiz"""
    user_id = update.effective_user.id
    
    # If quiz_id was provided with command, start that quiz directly
    if context.args:
        try:
            quiz_id = int(context.args[0])
            await start_saved_quiz(update, context, quiz_id)
            return
        except ValueError:
            await update.message.reply_text("מזהה מבחן לא תקין.")
            return
    
    # Otherwise, show list of quizzes
    with get_db_connection() as conn:
        c = conn.cursor()

        try:
            c.execute('''
                SELECT quiz_id, title, created_at
                FROM saved_quizzes
                WHERE user_id = ?
                ORDER BY created_at DESC
                LIMIT ?
            ''', (user_id, MAX_SAVED_QUIZZES_DISPLAY))
            quizzes = c.fetchall()

            if not quizzes:
                await update.message.reply_text("אין מבחנים שמורים!")
                return

            message = "בחר מבחן על ידי שליחת המספר המתאים:\n\n"
            for quiz_id, title, created_at in quizzes:
                message += f"{quiz_id}. {title} ({created_at})\n"

            context.user_data['waiting_for_quiz_selection'] = True
            await update.message.reply_text(message)

        except Exception as e:
            logger.error(f"Error listing quizzes: {str(e)}")
            await update.message.reply_text("❌ שגיאה בטעינת המבחנים. אנא נסה שוב.")

async def start_saved_quiz(update: Update, context: ContextTypes.DEFAULT_TYPE, quiz_id: int) -> None:
    """Start a saved quiz after ID selection, with caching and optimized error handling."""
    user_id = update.effective_user.id
    
    try:
        # First check if we have access permission
        async def check_access():
            with get_db_connection() as conn:
                c = conn.cursor()
                c.execute('''
                    SELECT 1 FROM saved_quizzes 
                    WHERE quiz_id = ? AND user_id = ?
                ''', (quiz_id, user_id))
                return bool(c.fetchone())
        
        has_access = await resilient_db_operation(check_access, fallback=lambda: False)
        
        if not has_access:
            await update.message.reply_text("❌ מבחן לא נמצא או שאין לך הרשאה להשתמש בו.")
            return
        
        # Get questions with caching
        questions = await db_get_quiz_with_caching(quiz_id)
        
        if not questions:
            await update.message.reply_text("❌ לא נמצאו שאלות במבחן זה.")
            return
        
        # Randomize options for each question on every play
        randomized_questions = []
        for question in questions:
            randomized_questions.append(randomize_question_options(question))

        # Update play count (in a separate operation to avoid blocking)
        async def update_play_count():
            async with get_async_db_connection() as conn:
                try:
                    await conn.execute('''
                        UPDATE saved_quizzes
                        SET times_played = times_played + 1
                        WHERE quiz_id = ? AND user_id = ?
                    ''', (quiz_id, user_id))
                    await conn.commit()
                    # Invalidate cache for this quiz's metadata
                    query_cache.invalidate(f"user_{user_id}_quizzes")
                except Exception as e:
                    logger.error(f"Error updating play count: {e}")
        
        # Run the update in the background
        asyncio.create_task(update_play_count())

        # Set up quiz in context
        context.user_data["randomized_questions"] = randomized_questions
        context.user_data["original_questions"] = randomized_questions.copy()  # Save a copy of the original questions
        context.user_data["current_question_index"] = 0
        context.user_data["correct_answers"] = 0
        context.user_data[STATE_QUIZ_ACTIVE] = True
        context.user_data["is_saved_quiz"] = True
        context.user_data["current_quiz_id"] = quiz_id

        # Start the quiz
        await send_question(update, context)

    except Exception as e:
        logger.error(f"Error starting saved quiz: {e}")
        await update.message.reply_text("❌ שגיאה בהפעלת המבחן. אנא נסה שוב.")

async def precheckout_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle the pre-checkout callback"""
    query = update.pre_checkout_query
    
    # Always approve the payment at this stage
    await query.answer(ok=True)

async def successful_payment_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle successful payments"""
    payment = update.message.successful_payment
    user_id = update.effective_user.id
    
    # Add one month of premium
    with get_db_connection() as conn:
        c = conn.cursor()
        
        try:
            # Calculate subscription end date (1 month from now)
            one_month_later = (datetime.now() + timedelta(days=30)).date()
            
            c.execute('''
                INSERT OR REPLACE INTO subscriptions (user_id, subscribed_until)
                VALUES (?, ?)
            ''', (user_id, one_month_later.isoformat()))
            
            conn.commit()
            
            await update.message.reply_text(
                "✨ תודה על הרכישה!\n"
                f"המנוי שלך פעיל עד {one_month_later.strftime('%d/%m/%Y')}\n"
                f"כעת יש לך {PREMIUM_USES_PER_DAY} ניסיונות יומיים."
            )
            
        except Exception as e:
            logger.error(f"Error processing successful payment: {e}")
            await update.message.reply_text("❌ אירעה שגיאה בעיבוד התשלום. אנא פנה לתמיכה.")

# Add support command handler
async def support_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Provide support contact information"""
    await update.message.reply_text(
        "📞 *תמיכת סיכום.AI*\n\n"
        "לכל שאלה או בעיה, ניתן ליצור קשר עם התמיכה:\n"
        f"{SUPPORT_CONTACT}\n\n"
        "אנא ציין:\n"
        "• את מהות הפנייה\n"
        "• מזהה משתמש שלך\n"
        "• צילום מסך של הבעיה (אם רלוונטי)\n\n"
        "זמן תגובה ממוצע: עד 24 שעות",
        parse_mode="Markdown"
    )

# Add cleanup of temp files
def cleanup_temp_files():
    """Clean up old temporary files."""
    for file in os.listdir(TEMP_DIR):
        file_path = os.path.join(TEMP_DIR, file)
        # Remove files older than 1 hour
        if os.path.getmtime(file_path) < time.time() - 3600:
            try:
                os.remove(file_path)
            except Exception as e:
                logger.error(f"Error cleaning up temp file {file}: {e}")

# Add after other constants
# Conversation States
WAITING_FEEDBACK = 1
WAITING_QUIZ_TITLE = 2
WAITING_QUIZ_SELECTION = 3

async def share_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Share a quiz with other users"""
    # If quiz_id wasn't provided with command, wait for it
    if not context.args:
        context.user_data['waiting_for_share_id'] = True
        await update.message.reply_text("אנא ציין את מזהה המבחן שברצונך לשתף.")
        return
        
    try:
        quiz_id = int(context.args[0])
        await share_quiz(update, context, quiz_id)
    except ValueError:
        await update.message.reply_text("מזהה מבחן לא תקין.")

async def share_quiz(update: Update, context: ContextTypes.DEFAULT_TYPE, quiz_id: int) -> None:
    """Share a saved quiz with others via a shareable link."""
    user_id = update.effective_user.id
    
    with get_db_connection() as conn:
        c = conn.cursor()
        
        # First check if quiz exists
        c.execute('SELECT 1 FROM saved_quizzes WHERE quiz_id = ?', (quiz_id,))
        if not c.fetchone():
            # Handle response based on whether it's a callback query or direct message
            if update.callback_query:
                await update.callback_query.edit_message_text("❌ מבחן זה לא נמצא.")
            else:
                await update.message.reply_text("❌ מבחן זה לא נמצא.")
            return
            
    share_link = f"https://t.me/{context.bot.username}?start=quiz_{quiz_id}"
    
    # Handle response based on whether it's a callback query or direct message
    if update.callback_query:
        await update.callback_query.edit_message_text(
            f"🔗 הנה הקישור לשיתוף המבחן:\n{share_link}",
            disable_web_page_preview=True
        )
    else:
        await update.message.reply_text(
            f"🔗 קישור לשיתוף המבחן:\n{share_link}",
            disable_web_page_preview=True
        )

# Add scheduled cleanup
from apscheduler.schedulers.asyncio import AsyncIOScheduler

# Add proper rate limiting decorator
from asyncio import Lock

class RateLimiter:
    def __init__(self):
        self.locks = {}
        
    async def acquire(self, user_id: int) -> bool:
        if user_id not in self.locks:
            self.locks[user_id] = Lock()
        
        if not self.locks[user_id].locked():
            await self.locks[user_id].acquire()
            return True
        return False

rate_limiter = RateLimiter()

def rate_limit(func):
    @wraps(func)
    async def wrapper(update: Update, context: ContextTypes.DEFAULT_TYPE, *args, **kwargs):
        user_id = update.effective_user.id
        if not await rate_limiter.acquire(user_id):
            await update.message.reply_text("נא להמתין מעט לפני שליחת בקשה נוספת.")
            return
        try:
            return await func(update, context, *args, **kwargs)
        finally:
            if user_id in rate_limiter.locks:
                rate_limiter.locks[user_id].release()
    return wrapper

# Removed duplicate optimize_db function

# Update the cleanup_old_sessions function to not require context
async def cleanup_old_sessions():
    """Cleanup user sessions older than 1 hour"""
    # Since we can't access context directly, we'll just clean temp files
    cleanup_temp_files()

async def cancel_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Cancel current operation."""
    # Clear any waiting states
    context.user_data.clear()
    
    await update.message.reply_text(
        "✅ הפעולה בוטלה.\n"
        "אתה יכול להתחיל מחדש עם אחת מהפקודות הזמינות:\n"
        "/start, /help, /play, /list, /save, /subscribe"
    )

async def handle_error(update: Update, context: ContextTypes.DEFAULT_TYPE, message: str = None):
    """Generic error handler with user notification."""
    try:
        user_id = update.effective_user.id if update and update.effective_user else None
        error_msg = f"Error: {context.error}" if context and context.error else "Unknown error"
        
        logger.error(f"Error for user {user_id}: {error_msg}")
        
        if user_id:
            error_text = message or "❌ אירעה שגיאה. אנא נסה שוב או פנה לתמיכה."
            try:
                if update.callback_query:
                    await update.callback_query.edit_message_text(text=error_text)
                else:
                    await context.bot.send_message(chat_id=user_id, text=error_text)
            except Exception as e:
                logger.error(f"Failed to send error message: {e}")
    except Exception as e:
        logger.error(f"Error in error handler: {e}")

async def subscribe_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle subscription purchases using Telegram Stars"""
    user_id = update.effective_user.id
    
    if await is_user_premium(user_id):
        await update.message.reply_text(
            "✨ אתה כבר מנוי פרימיום!\n"
            f"הגבלת השימוש היומית שלך היא {PREMIUM_USES_PER_DAY} ניסיונות."
        )
        return

    # Create the invoice for Telegram Stars payment
    await context.bot.send_invoice(
        chat_id=update.effective_chat.id,
        title="מנוי פרימיום סיכום.AI",
        description=(
            f"מנוי חודשי לסיכום.AI עם {PREMIUM_USES_PER_DAY} ניסיונות ביום"  # Updated to use constant
        ),
        payload=f"premium_sub_{user_id}",
        provider_token="",  # Empty for Telegram Stars
        currency="XTR",    # XTR is the currency code for Telegram Stars
        prices=[
            LabeledPrice(
                label="מנוי חודשי", 
                amount=100  # 100 Stars
            )
        ],
        need_name=False,
        need_phone_number=False,
        need_email=False,
        need_shipping_address=False,
        is_flexible=False,
        start_parameter="subscription"
    )

async def handle_text_input(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Handle text input for quiz selection and subscription"""
    if not update or not update.message:
        logger.error("Received update without message")
        return
        
    if not context or not context.user_data:
        await update.message.reply_text(
            "אנא התחל מחדש עם הפקודה /start"
        )
        return

    if context.user_data.get('waiting_for_quiz_selection'):
        try:
            quiz_id = int(update.message.text)
            context.user_data['waiting_for_quiz_selection'] = False
            await start_saved_quiz(update, context, quiz_id)
        except ValueError:
            await update.message.reply_text("אנא שלח מספר מבחן תקין.")
    elif context.user_data.get('waiting_for_quiz_title'):
        await handle_quiz_title(update, context)
    elif context.user_data.get('waiting_for_share_id'):
        try:
            quiz_id = int(update.message.text)
            context.user_data['waiting_for_share_id'] = False
            await share_quiz(update, context, quiz_id)
        except ValueError:
            await update.message.reply_text("אנא שלח מספר מבחן תקין.")
    else:
        await update.message.reply_text(
            "אנא השתמש באחת מהפקודות הזמינות:\n"
            "/start, /help, /play, /list, /save, /subscribe"
        )

# Add helper functions for subscription management
async def is_user_premium(user_id: int) -> bool:
    """Check if user has active premium subscription"""
    with get_db_connection() as conn:
        c = conn.cursor()
        
        try:
            c.execute('''
                SELECT subscribed_until FROM subscriptions 
                WHERE user_id = ? AND subscribed_until >= date('now')
            ''', (user_id,))
            result = c.fetchone()
            return bool(result)
        finally:
            conn.close()

async def get_user_daily_uses(user_id: int) -> int:
    """Get number of times user has used the service today."""
    with get_db_connection() as conn:
        c = conn.cursor()
        today = date.today().isoformat()
        
        c.execute('''
            SELECT attempts 
            FROM daily_usage 
            WHERE user_id = ? AND date = ?
        ''', (user_id, today))
        result = c.fetchone()
        return result[0] if result else 0

async def increment_user_daily_uses(user_id: int) -> None:
    """Increment user's daily usage count."""
    with get_db_connection() as conn:
        c = conn.cursor()
        today = date.today().isoformat()
        
        # First try to update existing record
        c.execute('''
            INSERT INTO daily_usage (user_id, date, attempts)
            VALUES (?, ?, 1)
            ON CONFLICT(user_id, date) 
            DO UPDATE SET attempts = attempts + 1
        ''', (user_id, today))
        
        conn.commit()

async def grant_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Secret command to grant premium access. Admin only."""
    user_id = update.effective_user.id
    
    # Check if user is admin
    if user_id not in ADMIN_IDS:
        # Don't reveal command exists
        await update.message.reply_text("פקודה לא קיימת.")
        return
        
    # Check if user ID was provided
    if not context.args:
        await update.message.reply_text("אנא ציין מזהה משתמש.")
        return
        
    try:
        target_user_id = int(context.args[0])
        with get_db_connection() as conn:
            c = conn.cursor()
            # Grant 30 days of premium
            one_month_later = (datetime.now() + timedelta(days=30)).date()
            c.execute('''
                INSERT OR REPLACE INTO subscriptions (user_id, subscribed_until)
                VALUES (?, ?)
            ''', (target_user_id, one_month_later.isoformat()))
            conn.commit()
            
        await update.message.reply_text(f"✅ הוענקה גישת פרימיום למשתמש {target_user_id}")
        logger.info(f"Admin {user_id} granted premium to user {target_user_id}")
        
    except ValueError:
        await update.message.reply_text("מזהה משתמש לא תקין.")
    except Exception as e:
        logger.error(f"Error granting premium: {e}")
        await update.message.reply_text("❌ אירעה שגיאה.")

def validate_quiz_title(title: str) -> bool:
    """Validate quiz title."""
    if not title or len(title) > 100:
        return False
    return True

async def check_time_sync(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Admin command to check time synchronization."""
    if update.effective_user.id not in ADMIN_IDS:
        return
        
    telegram_time = datetime.fromtimestamp(update.message.date.timestamp())
    local_time = datetime.now()
    time_diff = abs((local_time - telegram_time).total_seconds())
    
    await update.message.reply_text(
        f"Time check:\n"
        f"• Telegram time: {telegram_time}\n"
        f"• Server time: {local_time}\n"
        f"• Difference: {time_diff:.2f} seconds\n\n"
        f"{'⚠️ WARNING: Time difference exceeds 10 seconds!' if time_diff > 10 else '✓ Time synchronized correctly.'}"
    )

def cleanup_db_connections():
    """Clean up idle database connections."""
    with _pool_lock:
        current_time = time.time()
        idle_threshold = 300  # 5 minutes
        
        for thread_id in list(_connection_pool.keys()):
            # Check if connection is old and unused
            if current_time - _last_used.get(thread_id, 0) > idle_threshold:
                try:
                    conn = _connection_pool.pop(thread_id)
                    conn.close()
                    _last_used.pop(thread_id, None)
                    logger.debug(f"Closed idle connection for thread {thread_id}")
                except Exception as e:
                    logger.error(f"Error closing connection: {e}")

def repair_json(text):
    """Advanced JSON repair function."""
    import re
    
    # Remove any escaped unicode that might cause issues
    text = re.sub(r'\\u[0-9a-fA-F]{4}', ' ', text)
    
    # Fix common JSON issues
    text = re.sub(r',\s*\]', ']', text)  # Remove trailing commas in arrays
    text = re.sub(r',\s*\}', '}', text)  # Remove trailing commas in objects
    text = re.sub(r'",\s*(\s+)\]', '"\n\\1]', text)  # Fix trailing commas in multiline arrays
    
    # Fix broken quotes
    text = re.sub(r'(?<!")"(?!")', '\\"', text)
    
    # Remove control characters
    text = re.sub(r'[\x00-\x1F\x7F]', '', text)
    
    # Fix specific issue with options array missing closing bracket
    text = re.sub(r'(".+?"),?\s*"correct_option_index"', r'\1], "correct_option_index"', text)
    
    # Balance brackets if needed
    open_brackets = text.count('[') 
    close_brackets = text.count(']')
    if open_brackets > close_brackets:
        text += ']' * (open_brackets - close_brackets)
    
    # Fix missing commas between options
    text = re.sub(r'"\s*"', '", "', text)
    
    # Fix object structure: ensure each question has all required fields
    pattern = r'"question":[^}]*?"options":[^}]*?(?!"correct_option_index")(?=\s*\})'
    text = re.sub(pattern, lambda m: m.group(0) + ', "correct_option_index": 0, "explanation": "Auto-generated explanation"', text)
    
    # Ensure the whole text is a valid JSON array
    if not text.strip().startswith('['):
        text = '[' + text.strip()
    if not text.strip().endswith(']'):
        text = text.strip() + ']'
        
    logger.info(f"Repaired JSON: {text[:100]}...")
    return text

def main() -> None:
    """Start the bot with optimized settings."""
    # Print cool startup banner
    startup_banner = """
    ███████╗██╗██╗  ██╗██╗   ██╗███╗   ███╗      █████╗ ██╗
    ██╔════╝██║██║ ██╔╝██║   ██║████╗ ████║     ██╔══██╗██║
    ███████╗██║█████╔╝ ██║   ██║██╔████╔██║     ███████║██║
    ╚════██║██║██╔═██╗ ██║   ██║██║╚██╔╝██║     ██╔══██║██║
    ███████║██║██║  ██╗╚██████╔╝██║ ╚═╝ ██║     ██║  ██║██║
    ╚══════╝╚═╝╚═╝  ╚═╝ ╚═════╝ ╚═╝     ╚═╝     ╚═╝  ╚═╝╚═╝
    """
    print(startup_banner)
    print("🤖 Sikum.AI Bot is starting up...")
    print("📚 Ready to generate quizzes and help students learn!")
    print("🚀 Bot is now running! Press Ctrl+C to stop.")
    print("=" * 70)

    # Import async db functions from db_utils
    from db_utils import init_db as async_init_db, optimize_db as async_optimize_db, start_maintenance_task

    # Initialize database with optimized async settings
    # Create an event loop to run async initialization
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    
    try:
        # Run async database initialization
        loop.run_until_complete(async_init_db())
        loop.run_until_complete(async_optimize_db())
        
        # Start the database maintenance task
        maintenance_task = start_maintenance_task()
        
        logger.info("Database initialized with optimized concurrent settings")
    except Exception as e:
        logger.error(f"Error initializing database: {str(e)}")
        # Continue with regular initialization as fallback
        init_db()
        optimize_db()

    # Create application with optimized settings for minimizing timeouts
    application = (
        ApplicationBuilder()
        .token(BOT_TOKEN)
        .connection_pool_size(16)
        .pool_timeout(30.0)
        .read_timeout(5)         # Reduced from 7
        .write_timeout(5)         # Reduced from 7
        .connect_timeout(5)       # Added explicit connect timeout
        .get_updates_read_timeout(42)
        .build()
    )
    
    # Add handlers
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("stats", stats_command))
    application.add_handler(CommandHandler("play", play_command))
    application.add_handler(CommandHandler("list", list_saved_quizzes))
    application.add_handler(CommandHandler("save", save_quiz))
    application.add_handler(CommandHandler("share", share_command))
    application.add_handler(CommandHandler("cancel", cancel_command))
    application.add_handler(CommandHandler("subscribe", subscribe_command))
    application.add_handler(CommandHandler("support", support_command))
    application.add_handler(CommandHandler("grant", grant_command))
    application.add_handler(CallbackQueryHandler(handle_answer))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text_input))
    application.add_handler(PreCheckoutQueryHandler(precheckout_callback))
    application.add_handler(MessageHandler(filters.SUCCESSFUL_PAYMENT, successful_payment_callback))
    application.add_error_handler(error_handler)

    # Set up and start scheduler
    scheduler = AsyncIOScheduler()
    scheduler.add_job(cleanup_temp_files, 'interval', hours=1)
    scheduler.add_job(cleanup_old_sessions, 'interval', minutes=30)
    scheduler.add_job(cleanup_db_connections, 'interval', minutes=10)
    scheduler.start()

    # Add time check command
    application.add_handler(CommandHandler("timecheck", check_time_sync))

    # Start the bot
    print("Starting bot...")
    application.run_polling(drop_pending_updates=True)

if __name__ == '__main__':
    try:
        main()
    except KeyboardInterrupt:
        print("Received keyboard interrupt - shutting down")
    except Exception as e:
        print(f"Fatal error: {e}")
        traceback.print_exc()

# Add caching for frequently accessed data
from functools import lru_cache
import time

# Cache size and timeout settings
CACHE_SIZE = 100  # Number of quizzes to cache
CACHE_TTL = 300   # 5 minutes cache lifetime

# Cache quiz data to reduce database load
class QuizCache:
    def __init__(self, max_size=CACHE_SIZE, ttl=CACHE_TTL):
        self.cache = {}
        self.max_size = max_size
        self.ttl = ttl
        self.hits = 0
        self.misses = 0
    
    def get(self, quiz_id):
        """Get quiz from cache if available and not expired."""
        if quiz_id in self.cache:
            entry = self.cache[quiz_id]
            if time.time() - entry['timestamp'] < self.ttl:
                self.hits += 1
                return entry['data']
            # Expired entry
            del self.cache[quiz_id]
        self.misses += 1
        return None
    
    def set(self, quiz_id, data):
        """Store quiz in cache with timestamp."""
        # If cache is full, remove oldest entry
        if len(self.cache) >= self.max_size:
            oldest_id = min(self.cache.items(), key=lambda x: x[1]['timestamp'])[0]
            del self.cache[oldest_id]
        
        self.cache[quiz_id] = {
            'data': data,
            'timestamp': time.time()
        }
    
    def invalidate(self, quiz_id=None):
        """Invalidate specific quiz or entire cache."""
        if quiz_id is None:
            self.cache.clear()
        elif quiz_id in self.cache:
            del self.cache[quiz_id]
    
    def get_stats(self):
        """Return cache statistics."""
        total = self.hits + self.misses
        hit_rate = (self.hits / total) * 100 if total > 0 else 0
        return {
            'size': len(self.cache),
            'max_size': self.max_size,
            'hits': self.hits,
            'misses': self.misses,
            'hit_rate': f"{hit_rate:.1f}%"
        }

# Initialize cache
quiz_cache = QuizCache()

# Helper functions to work with cache
async def get_quiz_with_caching(quiz_id):
    """Get quiz with caching for better performance."""
    # Try cache first
    cached_quiz = quiz_cache.get(quiz_id)
    if cached_quiz:
        return cached_quiz
    
    # Cache miss, get from database
    async def db_operation():
        with get_db_connection() as conn:
            c = conn.cursor()
            c.execute('SELECT questions_json FROM saved_quizzes WHERE quiz_id = ?', (quiz_id,))
            result = c.fetchone()
            if result:
                data = json.loads(result[0])
                # Store in cache for future requests
                quiz_cache.set(quiz_id, data)
                return data
            return None
    
    # Use resilient operation with retries
    return await resilient_db_operation(
        db_operation,
        fallback=lambda: None
    )